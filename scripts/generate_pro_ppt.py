"""
生成专业比赛答辩 PPT — 世界杯冠军预测 Agent

输出: docs/presentation/世界杯冠军预测Agent_路演.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 配色 ──
DARK_BG  = RGBColor(0x0D, 0x1B, 0x2A)  # 深海蓝
CARD_BG  = RGBColor(0x14, 0x27, 0x3D)  # 卡片蓝
GOLD     = RGBColor(0xF0, 0xC0, 0x40)  # 金色
WHITE    = RGBColor(0xF0, 0xF0, 0xF0)  # 白色
GRAY     = RGBColor(0x90, 0x90, 0xA0)  # 灰色
GREEN    = RGBColor(0x2E, 0xCC, 0x71)  # 绿色
BLUE_ACC = RGBColor(0x34, 0x98, 0xDB)  # 蓝色强调
RED_ACC  = RGBColor(0xE7, 0x4C, 0x3C)  # 红色

ASSETS = os.path.join(os.path.dirname(__file__), "..", "docs", "assets")
OUTPUT = os.path.join(os.path.dirname(__file__), "..", "docs", "presentation",
                       "世界杯冠军预测Agent_路演.pptx")
os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# =====================================================================
# 辅助函数
# =====================================================================

def set_slide_bg(slide, color=DARK_BG):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gold_line(slide, left, top, width, height=Inches(0.03)):
    """添加金色分割线"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title="", lines=None, color=CARD_BG):
    """添加卡片组件"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = GRAY
    shape.line.width = Pt(0.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)

    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = GOLD

    if lines:
        for line in lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.space_after = Pt(3)

    return shape

def add_title(slide, text, left=Inches(0.8), top=Inches(0.3), size=Pt(32)):
    """添加页面标题"""
    txBox = slide.shapes.add_textbox(left, top, Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = True
    p.font.color.rgb = WHITE

def add_body_text(slide, lines, left=Inches(0.8), top=Inches(1.5), width=Inches(7)):
    """添加正文"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.startswith("## "):
            p.text = line[3:]
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = GOLD
            p.space_after = Pt(10)
        elif line.startswith("- "):
            p.text = "  ▸ " + line[2:]
            p.font.size = Pt(15)
            p.font.color.rgb = WHITE
            p.space_after = Pt(6)
        elif line.startswith("**") and line.endswith("**"):
            p.text = line[2:-2]
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = GOLD
            p.space_after = Pt(6)
        else:
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY

def add_image(slide, img_name, left, top, width, height=None):
    """添加素材图片"""
    img_path = os.path.join(ASSETS, img_name)
    if not os.path.exists(img_path):
        # 占位框
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height or Inches(3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = GRAY
        tf = shape.text_frame
        tf.paragraphs[0].text = f"[{img_name}]"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = GRAY
        return

    if height:
        slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        slide.shapes.add_picture(img_path, left, top, width)

def add_footer(slide, page_num):
    """添加页脚"""
    txBox = slide.shapes.add_textbox(Inches(10), Inches(7.1), Inches(3), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"世界杯冠军预测 Agent | {page_num}"
    p.font.size = Pt(9)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


# =====================================================================
# Slide 1: 封面
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# 顶部色条
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                prs.slide_width, Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

# 底部色条
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.44),
                                prs.slide_width, Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = GOLD; shape.line.fill.background()

# 装饰圆形
for cx, cy, r, alpha in [(Inches(3), Inches(2.5), Inches(1.8), 0.08),
                           (Inches(10.5), Inches(5), Inches(2.5), 0.05),
                           (Inches(1), Inches(5.5), Inches(1.2), 0.06)]:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r, r)
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    circle.fill.fore_color.brightness = 0.0

# 标题
txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(2.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "世界杯冠军预测 Agent"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "可解释 · 可视化 · 全流程自动化"
p2.font.size = Pt(22)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(16)

# 副标题
p3 = tf.add_paragraph()
p3.text = "基于五维度评分 + 泊松模型 + 蒙特卡洛模拟的智能预测系统"
p3.font.size = Pt(14)
p3.font.color.rgb = GRAY
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(30)

# 标签
tags = ["可解释预测", "可视化赛程树", "真实数据驱动", "工具链完整"]
tag_y = Inches(5.5)
for i, tag in enumerate(tags):
    tx = Inches(3.5 + i * 1.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx, tag_y, Inches(1.6), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x5C)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.paragraphs[0].text = tag
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = GOLD
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

add_footer(slide, "封面")


# =====================================================================
# Slide 2: 目录
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title(slide, "📋 目录")
add_gold_line(slide, Inches(0.8), Inches(1.0), Inches(4))

toc_items = [
    ("01", "项目背景与痛点", "世界杯预测市场现状与三大核心痛点"),
    ("02", "解决方案总览", "端到端可解释预测智能体架构"),
    ("03", "可解释预测引擎", "五维度评分 + 泊松 + 蒙特卡洛"),
    ("04", "系统架构", "三层架构与数据流设计"),
    ("05", "核心能力演示", "赛程树 · 评分雷达 · 因子归因"),
    ("06", "冠军预测结果", "蒙特卡洛 10000 次模拟 + Wilson CI"),
    ("07", "工具链与工程化", "Streamlit + FastAPI + Docker"),
    ("08", "竞品对比", "vs 博彩赔率 / 媒体专家 / AI 黑盒"),
    ("09", "项目亮点总结", "可解释 · 可视化 · 真实数据 · 可部署"),
    ("10", "部署与扩展", "阿里云 + Docker + 迭代路线图"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.6 + i * 0.55)
    col = i // 5
    x = Inches(0.8 + col * 6.2)

    # 编号
    txBox = slide.shapes.add_textbox(x, y, Inches(0.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # 标题
    txBox = slide.shapes.add_textbox(x + Inches(0.5), y, Inches(5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # 描述
    txBox = slide.shapes.add_textbox(x + Inches(0.5), y + Inches(0.25), Inches(5), Inches(0.25))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY

add_footer(slide, "目录")


# =====================================================================
# Slide 3: 项目背景
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title(slide, "01  项目背景与痛点")
add_gold_line(slide, Inches(0.8), Inches(1.0), Inches(4))

cards = [
    ("🎯 黑盒预测", "多数 AI 预测只给结果\n无法解释\"为什么 A 队会赢\"\n缺乏可解释性 → 用户不信任"),
    ("📊 静态输出", "只给出最终冠军名字\n缺少从小组赛到决赛的\n完整动态晋级路径"),
    ("🔗 数据碎片化", "历史战绩、排名、球员数据\n分散在不同来源\n人工整合成本极高"),
]

for i, (title, body) in enumerate(cards):
    x = Inches(0.8 + i * 4.1)
    add_card(slide, x, Inches(1.5), Inches(3.7), Inches(2.5),
             title=title, lines=body.split("\n"))

# 底部数据
add_body_text(slide, [
    "**数据支撑**",
    "- 2022 世界杯全球观众 50 亿+ | 体育预测市场 $1.5B+",
    "- \"可解释 AI\" 被列为 Gartner 2024 十大战略技术趋势",
    "- 本方案: 针对性地解决三大痛点，打造可解释、可视化、全流程闭环",
], top=Inches(4.5))

add_footer(slide, "3/12")


# =====================================================================
# Slide 4: 解决方案总览
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title(slide, "02  解决方案总览")
add_gold_line(slide, Inches(0.8), Inches(1.0), Inches(4))

add_image(slide, "architecture.png", Inches(0.8), Inches(1.3), Inches(11.7), Inches(4.7))

# 底部三个要点
caps = [
    ("真实数据驱动", "6 个公开数据源\n44K+ 比赛 · 64K+ 排名"),
    ("蒙特卡洛冠军预测", "10000 次模拟\nWilson 95% 置信区间"),
    ("交互式可视化", "赛程树 · 因子瀑布图\n雷达对比 · 热力图"),
]
for i, (t, body) in enumerate(caps):
    add_card(slide, Inches(0.8 + i * 4.1), Inches(6.2), Inches(3.7), Inches(1.0),
             title=t, lines=body.split("\n"))

add_footer(slide, "4/12")


# =====================================================================
# Slide 5: 可解释预测引擎
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title(slide, "03  可解释预测引擎")
add_gold_line(slide, Inches(0.8), Inches(1.0), Inches(4))

add_image(slide, "prediction_flow.png", Inches(0.5), Inches(1.3), Inches(7.8), Inches(5.2))

# 右侧说明
add_body_text(slide, [
    "**四层可解释架构**",
    "- Layer 0: 五维度球队评分 (历史+排名+攻防+球员+状态)",
    "- Layer 1: 独立泊松单场预测 (λ → 胜/平/负概率)",
    "- Layer 2: 赛程推演 (小组赛积分 + 淘汰赛三阶段)",
    "- Layer 3: 蒙特卡洛聚合 (N次模拟 → 夺冠概率 + CI)",
    "",
    "**关键特性**",
    "- 每一层都是白盒，可独立验证",
    "- 固定随机种子 (42)，结果完全可复现",
    "- 三阶段淘汰赛: 90min → 加时×0.33 → 点球55%",
], left=Inches(8.8), top=Inches(1.3), width=Inches(4.2))

add_footer(slide, "5/12")


# =====================================================================
# Slide 6: 系统架构
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title(slide, "04  系统架构与数据流")
add_gold_line(slide, Inches(0.8), Inches(1.0), Inches(4))

add_image(slide, "data_flow.png", Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.8))

add_body_text(slide, [
    "**技术栈**: Streamlit (前端) + FastAPI (后端) + NumPy/SciPy (模型) + Pandas (数据) + Docker (部署)",
    "- 数据层对接 Kaggle + FootballData 真实数据源，全链路可追溯",
    "- 服务层包含评分、预测、模拟、归因四大引擎，模块解耦",
    "- 交互层提供 5 页面 Streamlit 应用 + REST API",
], top=Inches(6.2))

add_footer(slide, "6/12")


# =====================================================================
# Slide 7: 核心能力演示
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title(slide, "05  核心能力演示")
add_gold_line(slide, Inches(0.8), Inches(1.0), Inches(4))

# 左: 赛程树
add_image(slide, "bracket_tree.png", Inches(0.3), Inches(1.3), Inches(7.2), Inches(3.0))

# 右: 雷达图
add_image(slide, "radar_chart.png", Inches(7.8), Inches(1.3), Inches(5.0), Inches(5.0))

add_body_text(slide, [
    "- **赛程树**: 从 32 强到冠军的完整晋级路径，悬停交互查看详情",
    "- **评分雷达**: 五维度多队对比，直观展示优劣势",
    "- **因子瀑布图**: 每场预测的可解释归因分析",
], left=Inches(0.5), top=Inches(4.5), width=Inches(7))

add_footer(slide, "7/12")


# =====================================================================
# Slide 8: 冠军预测结果
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title(slide, "06  冠军预测结果")
add_gold_line(slide, Inches(0.8), Inches(1.0), Inches(4))

add_image(slide, "champion_chart.png", Inches(0.5), Inches(1.3), Inches(7.5), Inches(4.5))

add_body_text(slide, [
    "**Top-3 热门**",
    "- 🥇 阿根廷 21.5%   [CI: 20.7%, 22.3%]",
    "- 🥈 巴西    18.5%   [CI: 17.8%, 19.2%]",
    "- 🥉 法国    14.2%   [CI: 13.5%, 14.9%]",
    "",
    "**关键指标**",
    "- 蒙特卡洛 10,000 次独立模拟",
    "- Wilson Score 95% 置信区间",
    "- 可复现: 固定随机种子 42",
], left=Inches(8.5), top=Inches(1.3), width=Inches(4.5))

add_footer(slide, "8/12")


# =====================================================================
# Slide 9: 工具链与工程化
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title(slide, "07  工具链与工程化")
add_gold_line(slide, Inches(0.8), Inches(1.0), Inches(4))

caps = [
    ("🖥️  前端 (5页面)", "Streamlit + Plotly\n赛程树 · 雷达图 · 瀑布图\n交互式参数调节"),
    ("🔌  后端 API", "FastAPI + Uvicorn\nRESTful 5 端点\n自动 Swagger 文档"),
    ("🐳  部署", "Docker + docker-compose\n阿里云 ECS 一键部署\n月费 < 200 元"),
    ("📦  工程规范", "pyproject.toml 依赖管理\n单元测试 · CLI 脚本\nGit 版本控制"),
]

for i, (t, body) in enumerate(caps):
    x = Inches(0.8 + (i % 2) * 6.2)
    y = Inches(1.5 + (i // 2) * 2.8)
    add_card(slide, x, y, Inches(5.8), Inches(2.5), title=t, lines=body.split("\n"))

add_footer(slide, "9/12")


# =====================================================================
# Slide 10: 竞品对比
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title(slide, "08  竞品对比与差异化")
add_gold_line(slide, Inches(0.8), Inches(1.0), Inches(4))

competitors = [
    ("博彩网站赔率", "只给数字，不解释原因", "五维度因子归因 + 瀑布图", RED_ACC, GREEN),
    ("媒体专家预测", "主观判断，不可复现", "数据驱动 + 固定种子可复现", RED_ACC, GREEN),
    ("AI 黑盒预测", "无法解释推理过程", "白盒架构，每层可独立验证", RED_ACC, GREEN),
]

for i, (name, pain, advantage, color1, color2) in enumerate(competitors):
    y = Inches(1.6 + i * 1.7)
    add_card(slide, Inches(0.8), y, Inches(5.5), Inches(1.4),
             title=name, lines=[f"痛点: {pain}", f"我们的优势: {advantage}"])

add_body_text(slide, [
    "**核心壁垒**",
    "- 可解释性: 不是给一个答案，而是让用户理解答案",
    "- 可视化: 交互式赛程树，竞品无同类产品",
    "- 完整性: 数据 → 预测 → 可视化 → 报告，全流程闭环",
    "- 工具链: Qoder + QoderWork 全流程开发",
], left=Inches(6.8), top=Inches(1.6), width=Inches(6))

add_footer(slide, "10/12")


# =====================================================================
# Slide 11: 项目亮点总结
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title(slide, "09  项目亮点总结")
add_gold_line(slide, Inches(0.8), Inches(1.0), Inches(4))

highlights = [
    ("🔮", "可解释预测", "每场预测拆解为 5 维度因子归因\n用户理解\"为什么赢\"，不是黑盒概率"),
    ("🌳", "可视化赛程树", "从 32 强到冠军，完整晋级路径\n悬停交互、缩放、路径高亮"),
    ("📊", "真实数据驱动", "6 个 Kaggle/GitHub 公开数据源\n全链路可追溯，支持答辩审计"),
    ("⚡", "完整工具链", "pip install → streamlit run\nDocker 一键部署，工业级交付"),
    ("🎯", "蒙特卡洛可靠性", "10000 次模拟 + Wilson 95% CI\n固定种子，结果稳定可复现"),
    ("🏆", "比赛级输出", "PPT · 报告 · API · 演示视频\n可直接用于路演答辩"),
]

for i, (icon, title, body) in enumerate(highlights):
    x = Inches(0.8 + (i % 3) * 4.1)
    y = Inches(1.5 + (i // 3) * 2.8)
    add_card(slide, x, y, Inches(3.7), Inches(2.5),
             title=f"{icon}  {title}", lines=body.split("\n"))

add_footer(slide, "11/12")


# =====================================================================
# Slide 12: 部署与扩展
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_title(slide, "10  部署与扩展")
add_gold_line(slide, Inches(0.8), Inches(1.0), Inches(4))

# 当前部署
add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5),
         title="🚀 当前部署 (MVP)",
         lines=[
             "阿里云 ECS (2vCPU 4GB)",
             "Docker + docker-compose 一键启动",
             "Streamlit :8501 + FastAPI :8000",
             "月费 < 200 元",
         ])

# 生产部署
add_card(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(2.5),
         title="☁️ 生产部署 (规划)",
         lines=[
             "ACK 容器服务 + SLB 负载均衡",
             "OSS 数据存储 + SLS 日志监控",
             "CDN 加速 + HTTPS",
             "水平扩展 + 自动伸缩",
         ])

# 迭代路线
add_card(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.0),
         title="📅 迭代路线图",
         lines=[
             "v1.0 MVP: 核心预测 + Streamlit + Docker (已完成)",
             "v1.1 → 实时数据 + XGBoost 集成 + PPT 导出 (2周)",
             "v1.2 → 情景分析引擎 + 球员级分析 + 模拟动画 (3周)",
             "v2.0 → 多赛事支持 (欧洲杯/欧冠) + SaaS 化 + LLM Agent (1月)",
         ])

add_footer(slide, "12/12")


# =====================================================================
# Slide 13: 结尾
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# 装饰圆
for cx, cy, r in [(Inches(2), Inches(3), Inches(2.5)),
                    (Inches(11), Inches(4), Inches(2))]:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r, r)
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    circle.fill.fore_color.brightness = 0.0

txBox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(3))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "我们不仅预测谁会赢得世界杯"
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "更让每个人理解为什么"
p2.font.size = Pt(36)
p2.font.bold = True
p2.font.color.rgb = GOLD
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(16)

p3 = tf.add_paragraph()
p3.text = "🏆  可解释 · 可视化 · 全流程自动化"
p3.font.size = Pt(18)
p3.font.color.rgb = GRAY
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(30)

p4 = tf.add_paragraph()
p4.text = "感谢评委 | 欢迎提问"
p4.font.size = Pt(16)
p4.font.color.rgb = GRAY
p4.alignment = PP_ALIGN.CENTER
p4.space_before = Pt(40)

add_footer(slide, "谢谢")


# =====================================================================
# 保存
# =====================================================================
prs.save(OUTPUT)
print(f"\n{'='*60}")
print(f"  专业 PPT 已生成!")
print(f"{'='*60}")
print(f"  文件: {OUTPUT}")
print(f"  页数: {len(prs.slides)} 页")
print(f"  格式: .pptx (可用 PowerPoint / WPS / Google Slides 打开)")
print(f"{'='*60}")
