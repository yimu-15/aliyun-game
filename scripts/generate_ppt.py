"""
生成可编辑的答辩 PPT

用法:
    python scripts/generate_ppt.py
输出:
    docs/presentation/世界杯冠军预测Agent_答辩.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── 颜色方案 ──
BLUE = RGBColor(0x1F, 0x77, 0xB4)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)


def add_slide(title_text, body_lines=None, subtitle=None):
    """添加一页幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 顶部色带
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK

    # 分割线
    line = slide.shapes.add_shape(
        1, Inches(0.8), Inches(1.2), Inches(4), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    # 副标题
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY

    # 正文
    if body_lines:
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(5.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            # 标题行
            if line.startswith("## "):
                p.text = line[3:]
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = BLUE
                p.space_after = Pt(6)
            elif line.startswith("- "):
                p.text = "  • " + line[2:]
                p.font.size = Pt(16)
                p.font.color.rgb = DARK
                p.space_after = Pt(4)
            elif line.startswith("|"):
                p.text = line
                p.font.size = Pt(13)
                p.font.color.rgb = DARK
                p.font.name = "Consolas"
            elif line.startswith(">"):
                p.text = "  " + line[1:]
                p.font.size = Pt(15)
                p.font.italic = True
                p.font.color.rgb = GRAY
            elif line.strip():
                p.text = line
                p.font.size = Pt(16)
                p.font.color.rgb = DARK
                p.space_after = Pt(4)

    # 页脚
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "世界杯冠军预测 Agent | 可解释 · 可视化 · 全流程"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

    return slide


# =====================================================================
# Slide 1: 封面
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
# 背景色块
bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "🏆 世界杯冠军预测 Agent"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = GOLD
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "可解释 · 可视化 · 全流程自动化"
p2.font.size = Pt(24)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(20)

p3 = tf.add_paragraph()
p3.text = "基于五维度评分 + 泊松模型 + 蒙特卡洛模拟的智能预测系统"
p3.font.size = Pt(16)
p3.font.color.rgb = GRAY
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(40)

# 页脚
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.5))
tf2 = txBox2.text_frame
p = tf2.paragraphs[0]
p.text = "比赛路演答辩 | 2026"
p.font.size = Pt(14)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

# =====================================================================
# Slide 2-12: 内容页
# =====================================================================

add_slide("项目背景与痛点", [
    "## 世界杯预测的三大痛点",
    "",
    "| 痛点 | 表现 | 本方案 |",
    "|------|------|--------|",
    "| 黑盒预测 | 多数AI只给结果，无法解释原因 | SHAP + 因子消融 |",
    "| 静态输出 | 只给冠军，缺少动态晋级路径 | 交互式赛程树 |",
    "| 数据碎片化 | 历史战绩、排名、球员数据分散 | 统一数据底座 |",
    "",
    "- 2022 世界杯全球观众 50 亿+",
    "- 体育预测市场规模 $1.5B+",
    "- '可解释 AI' 是 Gartner 2024 十大战略技术趋势",
], "如何让冠军预测既准确又可解释？")

add_slide("解决方案总览", [
    "## 端到端的可解释预测智能体",
    "",
    "     数据采集层          预测引擎层          可视化展示层",
    "  ┌──────────────┐  ┌──────────────┐  ┌──────────────┐",
    "  │ Kaggle+FIFA   │  │ 5维评分+泊松  │  │ Streamlit    │",
    "  │ 40支球队数据  │→│ +蒙特卡洛     │→│ +Plotly      │",
    "  └──────────────┘  └──────────────┘  └──────────────┘",
    "",
    "- 真实数据驱动 (6 个 Kaggle/GitHub 公开数据源)",
    "- 五维度球队评分 (历史 + 排名 + 攻防 + 球员 + 状态)",
    "- 蒙特卡洛冠军预测 (10000 次模拟, 95% 置信区间)",
    "- 交互式可视化 (赛程树 + 因子瀑布图)",
], "三层架构: 数据层 → 引擎层 → 展示层")

add_slide("核心亮点一: 可解释预测", [
    "## 不仅告诉你'谁会赢'，更解释'为什么赢'",
    "",
    "  方法: 单项特征消融法 (Single Feature Ablation)",
    "",
    "  例: 巴西 vs 阿根廷 → 巴西胜率 62%",
    "",
    "  | 胜负因子    | 巴西  | 阿根廷 | 贡献度 |",
    "  |------------|-------|--------|--------|",
    "  | 近期状态    | 0.85  | 0.55   | +12% → |",
    "  | FIFA排名    | 0.92  | 0.78   | +6%  → |",
    "  | 历史交锋    | 0.60  | 0.40   | +5%  → |",
    "  | 攻防效率    | 0.72  | 0.68   | +3%  → |",
    "  | 球员班底    | 0.84  | 0.79   | +2%  → |",
    "",
    "- 瀑布图直观展示因子贡献",
    "- 每个因子可追溯到原始数据",
    "- 支持用户自定义权重",
], "解释机制: 逐因子消融 → 贡献度排序")

add_slide("核心亮点二: 可视化赛程树", [
    "## 从小组赛到冠军，完整动态晋级路径",
    "",
    "                    ┌──────────┐",
    "                    │ 🏆 冠军  │",
    "                    │ 巴西 58% │",
    "                    └────┬─────┘",
    "           ┌─────────────┴─────────────┐",
    "      ┌────┴────┐                 ┌────┴────┐",
    "      │ 半决赛1  │                 │ 半决赛2  │",
    "      │巴vs法 62%│                 │阿vs英 55%│",
    "      └────┬────┘                 └────┬────┘",
    "",
    "- 悬停节点 → 显示球队详情 + 晋级概率",
    "- 点击节点 → 高亮完整晋级路径",
    "- 支持缩放/平移/导出",
    "- 实线 = 已确定 | 虚线 = 预测",
], "交互能力: 悬停 · 点击 · 缩放 · 路径追踪")

add_slide("系统架构", [
    "## 三层架构，模块解耦",
    "",
    "  ┌─────────────────────────────────────────┐",
    "  │           用户交互层                      │",
    "  │  Streamlit 前端 (5页面) │ FastAPI (REST)  │",
    "  ├─────────────────────────────────────────┤",
    "  │           应用服务层                      │",
    "  │  ┌──────────┐ ┌──────────┐ ┌─────────┐ │",
    "  │  │球队评分   │ │单场预测   │ │蒙特卡洛  │ │",
    "  │  │5维加权   │ │泊松模型   │ │N=10000  │ │",
    "  │  └──────────┘ └──────────┘ └─────────┘ │",
    "  │  ┌──────────┐ ┌──────────┐ ┌─────────┐ │",
    "  │  │淘汰赛模拟 │ │因子归因   │ │报告生成  │ │",
    "  │  │3阶段制胜 │ │特征消融   │ │MD/JSON  │ │",
    "  │  └──────────┘ └──────────┘ └─────────┘ │",
    "  ├─────────────────────────────────────────┤",
    "  │           数据层                          │",
    "  │  数据采集(KaggleHub) │ 清洗 │ CSV存储     │",
    "  └─────────────────────────────────────────┘",
], "技术栈: Streamlit + FastAPI + NumPy + Plotly")

add_slide("数据采集与管理", [
    "## 真实数据驱动，全链路可追溯",
    "",
    "| 数据源 | 内容 | 数据量 |",
    "|--------|------|--------|",
    "| Kaggle - Intl Football | 国家队比赛记录 | 44,000+ 场 |",
    "| Kaggle - FIFA World Cup | 全部22届世界杯 | 900+ 场 |",
    "| Kaggle - FIFA Rankings | 每月排名快照 | 1992-2024 |",
    "| Kaggle - FIFA 23 Players | 球员能力值 | 19,000+ 球员 |",
    "| GitHub - openfootball | 世界杯赛程 | 全部届次 |",
    "| eloratings.net | ELO 评分 | 实时 |",
    "",
    "- 每条数据可追溯到源 URL + 下载时间",
    "- 缺失数据透明标注 + 真实替代方案",
], "6 个公开数据源，全链路可追溯")

add_slide("预测模型设计", [
    "## 四层可解释预测架构",
    "",
    "- **Layer 0: 球队评分** = 历史×20% + 排名×30% + 攻防×20% + 球员×15% + 状态×15%",
    "- **Layer 1: 泊松预测** λ = 1.32 × 攻击力(A) × 防守力(B) → P(胜/平/负)",
    "- **Layer 2: 赛程推演** 小组赛(泊松采样) + 淘汰赛(90min→加时→点球)",
    "- **Layer 3: 蒙特卡洛** N=10000次 → 夺冠概率 + Wilson 95% CI",
    "",
    "  关键特性:",
    "- 每一层都是白盒，可独立验证",
    "- 固定随机种子(42)，结果完全可复现",
    "- 三阶段淘汰赛: 90分钟 → 加时×0.33 → 点球55%",
], "白盒模型: 每一层都可独立验证和解释")

add_slide("技术实现与工具链", [
    "## 完整工具链，工业级工程实践",
    "",
    "  worldcup-predictor/",
    "  ├── config/      YAML 配置 (数据源/模型参数)",
    "  ├── data/        数据目录 (raw/processed/metadata)",
    "  ├── models/      核心引擎 (rating/predictor/sim/explainer)",
    "  ├── backend/     FastAPI 服务",
    "  ├── app/         Streamlit 前端 (5页面)",
    "  ├── utils/       工具 (logger/loader/helpers)",
    "  ├── scripts/     CLI 脚本 (fetch/pipeline/demo)",
    "  └── tests/       单元测试",
    "",
    "- `pip install -e .` 一键安装",
    "- `python main.py --teams 32 --sim 10000` 命令行运行",
    "- `streamlit run app/main.py` 可视化界面",
    "- 所有配置参数化 (权重/种子/赛制可调节)",
], "标准 Python 包结构，可交付可运行")

add_slide("项目成果与指标", [
    "## MVP 交付成果",
    "",
    "| 维度 | 目标 | 达成 |",
    "|------|------|------|",
    "| 功能完整性 | 数据→预测→可视化→报告 | ✅ 全流程打通 |",
    "| 可解释性 | ≥5 个可解释因子 | ✅ 五维度因子归因 |",
    "| 数据真实性 | 6 个公开数据源 | ✅ 全部可追溯 |",
    "| 可视化 | 交互式赛程树 | ✅ Plotly 交互 |",
    "| 可运行性 | 一键启动 | ✅ streamlit run |",
    "| 可配置性 | 权重/种子/赛制可调 | ✅ 侧边栏参数 |",
    "",
    "- 核心模块: ~1200 行 Python",
    "- Streamlit 页面: 5 个",
    "- CLI 脚本: 4 个",
], "代码规模: 1200+ 行核心代码，5 个可视化页面")

add_slide("竞品对比与差异化", [
    "## 我们 vs 竞品",
    "",
    "| 竞品 | 痛点 | 我们的优势 |",
    "|------|------|-----------|",
    "| 博彩赔率 | 只给数字，不解释原因 | 五维度因子归因 + 瀑布图 |",
    "| 媒体专家 | 主观判断，不可复现 | 数据驱动 + 固定种子可复现 |",
    "| AI 预测 | 黑盒模型 | 白盒架构，每层可验证 |",
    "",
    "- **核心壁垒**: 可解释性 + 赛程树 + 全工具链",
    "- **差异化**: 不是给一个答案，而是让用户理解答案",
], "核心壁垒: 可解释性 + 可视化 + 完整性")

add_slide("总结与展望", [
    "## 可解释预测，让AI不再是黑盒",
    "",
    "  项目价值:",
    "- 技术价值: 证明了'可解释预测'在体育领域的可行性",
    "- 产品价值: 完整的端到端系统，可立即使用",
    "- 社会价值: 让普通球迷也能理解AI预测的逻辑",
    "",
    "  后续计划:",
    "  v1.1 → 实时数据接入 + XGBoost + PPT导出",
    "  v1.2 → 情景分析引擎 + 球员级分析 + 模拟动画",
    "  v2.0 → 多赛事支持 + 用户系统 + SaaS化",
    "",
    "> '我们不仅预测谁会赢得世界杯，更让每个人理解为什么。'",
], "一句话总结: 预测 + 解释 + 可视化 = 完整闭环")

# =====================================================================
# 保存
# =====================================================================
output_dir = os.path.join(os.path.dirname(__file__), "..", "docs", "presentation")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "世界杯冠军预测Agent_答辩.pptx")
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
print("可用 PowerPoint / WPS / Google Slides 打开编辑")
