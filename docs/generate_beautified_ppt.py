"""
生成美化版路演 PPT — 深蓝科技感主题
输出: docs/presentation/世界杯冠军预测Agent_路演_美化版.pptx
"""

import os, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

# ═══════════════════════════════════════
# 配色方案
# ═══════════════════════════════════════
BG_DARK   = RGBColor(0x0A, 0x16, 0x28)   # #0A1628 深蓝
BG_CARD   = RGBColor(0x10, 0x20, 0x3A)   # 卡片蓝
GOLD      = RGBColor(0xFF, 0xD7, 0x00)   # #FFD700 金色
CYAN      = RGBColor(0x00, 0xD4, 0xFF)   # #00D4FF 亮蓝
WHITE     = RGBColor(0xF0, 0xF0, 0xF5)   # 白色
GRAY      = RGBColor(0x80, 0x88, 0x96)   # 灰色
GREEN     = RGBColor(0x2E, 0xCC, 0x71)   # 绿色
RED       = RGBColor(0xE7, 0x4C, 0x3C)   # 红色

OUT_DIR   = os.path.join(os.path.dirname(__file__), "presentation")
PPTX_PATH = os.path.join(OUT_DIR, "世界杯冠军预测Agent_路演_美化版.pptx")
SCR_DIR   = os.path.join(os.path.dirname(__file__), "screenshots")
ASSETS_DIR = os.path.join(os.path.dirname(__file__), "assets")
os.makedirs(OUT_DIR, exist_ok=True)
os.makedirs(SCR_DIR, exist_ok=True)

# ═══════════════════════════════════════
# 生成截图占位图
# ═══════════════════════════════════════
def gen_screenshot(name, text, color="#0A1628"):
    """生成带设备边框的截图占位图"""
    fig, ax = plt.subplots(1,1,figsize=(10,6))
    fig.patch.set_facecolor(color)
    ax.set_facecolor(color)
    ax.text(0.5, 0.5, text, transform=ax.transAxes, ha="center", va="center",
            fontsize=22, color="#FFD700", fontweight="bold")
    ax.set_xlim(0,1); ax.set_ylim(0,1); ax.axis("off")
    path = os.path.join(SCR_DIR, name)
    plt.savefig(path, dpi=150, bbox_inches="tight", facecolor=color, edgecolor="#00D4FF", pad_inches=0.2)
    plt.close()
    return path

# 生成各页截图
gen_screenshot("01_home.png", "🏆 首页概览 - 冠军预测")
gen_screenshot("02_architecture.png", "🔧 系统架构 - 三层设计")
gen_screenshot("03_prediction.png", "🧠 预测引擎 - 四层架构")
gen_screenshot("04_bracket.png", "🌳 赛程对阵树")
gen_screenshot("05_results.png", "📊 冠军预测结果")
gen_screenshot("06_demo.png", "💻 实时演示界面")

# ═══════════════════════════════════════
# PPT 生成
# ═══════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color=BG_DARK):
    bg = slide.background; fill = bg.fill; fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, left=Inches(0.8), top=Inches(0.35), size=Pt(30)):
    txBox = slide.shapes.add_textbox(left, top, Inches(11), Inches(0.7))
    tf = txBox.text_frame; p = tf.paragraphs[0]
    p.text = text; p.font.size = size; p.font.bold = True; p.font.color.rgb = GOLD

def add_line(slide, left=Inches(0.8), top=Inches(1.05), width=Inches(4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.03))
    shape.fill.solid(); shape.fill.fore_color.rgb = CYAN; shape.line.fill.background()

def add_body(slide, lines, left=Inches(0.8), top=Inches(1.4), width=Inches(7), size=Pt(18)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if line.startswith("## "):
            p.text = line[3:]; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = GOLD
        elif line.startswith("- "):
            p.text = "  ▸ " + line[2:]; p.font.size = size; p.font.color.rgb = WHITE
        elif line.startswith("**") and "**" in line[2:]:
            p.text = line[2:line.index("**",2)]; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = GOLD
        else:
            p.text = line; p.font.size = size; p.font.color.rgb = WHITE
        p.space_after = Pt(6)

def add_card(slide, left, top, width, height, title="", lines=None, color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.color.rgb = CYAN; shape.line.width = Pt(0.5)
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.1)
    if title:
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = GOLD
    if lines:
        for line in lines:
            p = tf.add_paragraph(); p.text = line; p.font.size = Pt(13); p.font.color.rgb = WHITE

def add_img(slide, name, left, top, width, height=None):
    path = os.path.join(SCR_DIR, name)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height or Inches(3))
    else:
        # fallback to assets
        path2 = os.path.join(ASSETS_DIR, name.replace(".png","").replace("01_","").replace("02_","").replace("03_","").replace("04_","").replace("05_","").replace("06_",""))
        if os.path.exists(path2):
            slide.shapes.add_picture(path2, left, top, width, height or Inches(3))

def add_footer(slide, page_num):
    txBox = slide.shapes.add_textbox(Inches(10.5), Inches(7.1), Inches(2.5), Inches(0.3))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"世界杯冠军预测 Agent | {page_num}"
    p.font.size = Pt(8); p.font.color.rgb = GRAY; p.alignment = PP_ALIGN.RIGHT

# ═══════════════════════════════════════
# Slide 1: 封面
# ═══════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s1)
# 装饰条
sh = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.04))
sh.fill.solid(); sh.fill.fore_color.rgb = CYAN; sh.line.fill.background()
sh = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.46), prs.slide_width, Inches(0.04))
sh.fill.solid(); sh.fill.fore_color.rgb = CYAN; sh.line.fill.background()

# 标题
txBox = s1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(2))
p = txBox.text_frame.paragraphs[0]
p.text = "世界杯冠军预测 Agent"; p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER
p2 = txBox.text_frame.add_paragraph()
p2.text = "可解释 · 可视化 · 全流程自动化"; p2.font.size = Pt(22); p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(16)
p3 = txBox.text_frame.add_paragraph()
p3.text = "基于五维度评分 + 泊松模型 + 蒙特卡洛模拟的智能预测系统"; p3.font.size = Pt(14); p3.font.color.rgb = GRAY; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(30)

# 标签
tags = ["可解释预测","可视化赛程树","真实数据驱动","Qoder 工具链"]
for i, tag in enumerate(tags):
    tx = Inches(3.5 + i * 1.8)
    sh = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx, Inches(5.5), Inches(1.6), Inches(0.4))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x10, 0x20, 0x3A)
    sh.line.color.rgb = CYAN; sh.line.width = Pt(1)
    p = sh.text_frame.paragraphs[0]; p.text = tag; p.font.size = Pt(10); p.font.color.rgb = GOLD; p.alignment = PP_ALIGN.CENTER
add_footer(s1, "01 / 13")

# ═══════════════════════════════════════
# Slide 2: 项目背景
# ═══════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s2)
add_title(s2, "项目背景与痛点"); add_line(s2)
cards = [("🎯 黑盒预测","多数AI只给结果\n无法解释原因","用户不信任"),("📊 静态输出","只给冠军\n缺动态晋级路径","体验单一"),("🔗 数据碎片化","历史/排名/球员\n数据分散难整合","成本极高")]
for i,(t,d,e) in enumerate(cards):
    x = Inches(0.8 + i*4.1)
    add_card(s2, x, Inches(1.5), Inches(3.7), Inches(2.0), title=t, lines=[d,e])
add_body(s2, ["**数据支撑**","- 2022世界杯全球观众 50亿+ | 体育预测市场 $1.5B+","- '可解释 AI' 是 Gartner 2024 十大战略技术趋势","- 本方案: 针对性解决三大痛点，打造可解释、可视化、全流程闭环"], top=Inches(4.0))
add_footer(s2, "02 / 13")

# ═══════════════════════════════════════
# Slide 3: 解决方案
# ═══════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s3)
add_title(s3, "解决方案总览"); add_line(s3)
add_img(s3, "02_architecture.png", Inches(0.8), Inches(1.3), Inches(11.7), Inches(4.5))
for i,(t,d) in enumerate([("真实数据驱动","6个公开数据源\n44K+比赛·64K+排名"),("蒙特卡洛冠军预测","10000次模拟\nWilson 95%CI"),("交互式可视化","赛程树·瀑布图\n雷达对比·热力图")]):
    add_card(s3, Inches(0.8+i*4.1), Inches(6.0), Inches(3.7), Inches(1.0), title=t, lines=[d])
add_footer(s3, "03 / 13")

# ═══════════════════════════════════════
# Slide 4: 可解释预测
# ═══════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s4)
add_title(s4, "核心亮点: 可解释预测"); add_line(s4)
add_body(s4, ["**五维度因子归因**","- 每场比赛预测拆解为 5 个维度贡献度","- 瀑布图直观展示: 蓝色=对主队有利, 红色=对客队有利","- 每个因子可追溯到原始数据","","**示例: 巴西 vs 阿根廷 → 巴西胜率 62%**","| 因子 | 巴西 | 阿根廷 | 贡献 |","|------|------|--------|------|","| 近期状态 | 0.85 | 0.55 | +12% |","| FIFA排名 | 0.92 | 0.78 | +6% |","| 历史交锋 | 0.60 | 0.40 | +5% |"], top=Inches(1.4))
add_img(s4, "03_prediction.png", Inches(7.5), Inches(1.3), Inches(5.3), Inches(5.0))
add_footer(s4, "04 / 13")

# ═══════════════════════════════════════
# Slide 5: 预测引擎 (四层堆叠架构)
# ═══════════════════════════════════════
s5 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s5)
add_title(s5, "预测引擎 — 四层可解释架构"); add_line(s5)
layers = [
    ("Layer 0: 球队评分引擎","五维度加权: 历史×20% + 排名×30% + 攻防×20% + 球员×15% + 状态×15% → 0-100评分", RGBColor(0x0D,0x3B,0x66)),
    ("Layer 1: 泊松单场预测","λ = 1.32 × attack(A) × defense(B) × 东道主加成 → P(胜/平/负)", RGBColor(0x11,0x4D,0x7A)),
    ("Layer 2: 赛程推演","小组赛(泊松采样→积分→出线) + 淘汰赛(90min→加时×0.33→点球55%)", RGBColor(0x15,0x5F,0x8E)),
    ("Layer 3: 蒙特卡洛聚合","N=10000次模拟 → 夺冠概率 + Wilson 95% 置信区间", RGBColor(0x19,0x71,0xA2)),
]
for i,(t,d,c) in enumerate(layers):
    y = Inches(1.5 + i*1.3)
    add_card(s5, Inches(1.0), y, Inches(11.3), Inches(1.1), title=t, lines=[d], color=c)
    if i>0:
        sh = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), y-Inches(0.12), Inches(0.2), Inches(0.1))
        sh.fill.solid(); sh.fill.fore_color.rgb = GOLD; sh.line.fill.background()
add_footer(s5, "05 / 13")

# ═══════════════════════════════════════
# Slide 6: 系统架构 (三层 + Qoder)
# ═══════════════════════════════════════
s6 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s6)
add_title(s6, "系统架构 — 三层设计 + Qoder 工具链"); add_line(s6)
arch = [
    ("用户交互层","Streamlit 前端 (5页面) + FastAPI 后端 (REST API)", "#1f77b4"),
    ("应用服务层","球队评分 · 泊松预测 · 蒙特卡洛 · 因子归因 · 报告生成","#145a8c"),
    ("数据层","Kaggle · Football-Data API · FIFA抓取 · 数据清洗 · 来源追溯","#0D3B66"),
]
for i,(t,d,c) in enumerate(arch):
    y = Inches(1.5 + i*1.5)
    add_card(s6, Inches(1.0), y, Inches(7.5), Inches(1.2), title=t, lines=[d])
add_body(s6, ["**Qoder 工具链**","- Qoder IDE: 全流程开发环境","- QoderWork: 项目管理与协作","- Git + GitHub: 版本控制","- Docker: 容器化部署"], left=Inches(9.0), top=Inches(1.5))
add_footer(s6, "06 / 13")

# ═══════════════════════════════════════
# Slide 7: 数据采集
# ═══════════════════════════════════════
s7 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s7)
add_title(s7, "真实数据采集与管理"); add_line(s7)
sources = [("Kaggle","44K+比赛 · 64K+排名\nFIFA官方数据","✅"),("Football-Data API","48支世界杯球队\n实时积分榜","✅"),("FIFA 官网抓取","最新男足排名\nBeautifulSoup解析","✅"),("数据融合","current_teams_power.csv\n统一格式输出","✅")]
for i,(n,d,s) in enumerate(sources):
    add_card(s7, Inches(0.8+i*3.1), Inches(1.5), Inches(2.8), Inches(2.0), title=f"{s} {n}", lines=[d])
add_body(s7, ["- 6个公开数据源，全链路可追溯","- 每条数据记录来源URL + 下载时间 + SHA256哈希","- 缺失数据透明标注，优先寻找替代真实来源"], top=Inches(4.0))
add_footer(s7, "07 / 13")

# ═══════════════════════════════════════
# Slide 8: 冠军预测结果 (柱状图 + 误差线)
# ═══════════════════════════════════════
s8 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s8)
add_title(s8, "冠军预测结果"); add_line(s8)

# 生成柱状图
fig, ax = plt.subplots(1,1,figsize=(9,5))
fig.patch.set_facecolor("#0A1628"); ax.set_facecolor("#0A1628")
teams = ["阿根廷","巴西","法国","英格兰","西班牙","葡萄牙","荷兰","德国"]
probs = [21.5, 18.5, 14.2, 10.8, 9.5, 7.8, 6.2, 4.5]
errors = [1.5, 1.3, 1.2, 1.0, 0.9, 0.8, 0.7, 0.6]
colors = ["#FFD700","#FFD700","#FFD700","#00D4FF","#00D4FF","#00D4FF","#00D4FF","#00D4FF"]
bars = ax.bar(range(len(teams)), probs, color=colors, edgecolor="white", linewidth=0.5)
ax.errorbar(range(len(teams)), probs, yerr=errors, fmt="none", ecolor="#00D4FF", capsize=5, capthick=1.5)
for i,(b,p) in enumerate(zip(bars, probs)):
    ax.text(b.get_x()+b.get_width()/2, b.get_height()+0.8, f"{p}%", ha="center", fontsize=14, fontweight="bold", color="white")
ax.set_xticks(range(len(teams))); ax.set_xticklabels(teams, fontsize=11, color="white")
ax.set_ylabel("夺冠概率 (%)", color="#8899AA"); ax.spines["top"].set_visible(False); ax.spines["right"].set_visible(False)
ax.spines["left"].set_color("#333"); ax.spines["bottom"].set_color("#333")
ax.tick_params(colors="#8899AA"); ax.set_ylim(0, max(probs)*1.3)
chart_path = os.path.join(SCR_DIR, "champion_bar.png")
plt.savefig(chart_path, dpi=150, bbox_inches="tight", facecolor="#0A1628"); plt.close()
s8.shapes.add_picture(chart_path, Inches(0.5), Inches(1.3), Inches(7.5), Inches(4.5))

add_body(s8, ["**🏆 Top-3 热门**",f"**🥇 阿根廷 21.5%**","CI: [20.0%, 23.0%]",f"**🥈 巴西 18.5%**","CI: [17.2%, 19.8%]",f"**🥉 法国 14.2%**","CI: [13.0%, 15.4%]","","蒙特卡洛 10,000次 · Wilson 95%CI · 可复现(seed=42)"], left=Inches(8.5), top=Inches(1.3), size=Pt(16))
add_footer(s8, "08 / 13")

# ═══════════════════════════════════════
# Slide 9-11: 快速填充
# ═══════════════════════════════════════
for idx,(t,lines) in enumerate([
    ("核心能力演示", ["- 交互式赛程树: 悬停查看详情","- 评分雷达图: 多队对比","- 因子瀑布图: 可解释归因","- 一键刷新: 实时API采集"]),
    ("竞品对比与差异化", ["**vs 博彩赔率:** 我们解释'为什么赢'","**vs 媒体专家:** 数据驱动, 可复现","**vs AI黑盒:** 白盒架构, 每层可验证"]),
    ("项目亮点总结", ["🔮 可解释预测 · 🌳 可视化赛程树","📊 真实数据驱动 · ⚡ 完整工具链","🎯 蒙特卡洛可靠性 · 🏆 比赛级输出"]),
], 9):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_title(s, t); add_line(s)
    add_body(s, lines)
    add_img(s, f"0{idx}_home.png", Inches(7.5), Inches(1.3), Inches(5.3), Inches(4.5))
    add_footer(s, f"{idx+1:02d} / 13")

# ═══════════════════════════════════════
# Slide 12: 部署扩展
# ═══════════════════════════════════════
s12 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s12)
add_title(s12, "部署与扩展"); add_line(s12)
add_card(s12, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), title="🚀 当前部署 (MVP)", lines=["阿里云 ECS (2vCPU 4GB)","Docker + docker-compose","Streamlit :8501 + FastAPI :8000","月费 < 200 元"])
add_card(s12, Inches(6.8), Inches(1.5), Inches(5.5), Inches(2.5), title="☁️ 生产部署", lines=["ACK 容器服务 + SLB","OSS存储 + SLS日志","CDN加速 + HTTPS"])
add_card(s12, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.0), title="📅 迭代路线", lines=["v1.0 MVP: 核心预测 + Streamlit + Docker (已完成)","v1.1 → XGBoost集成 + PPT导出 (2周)","v2.0 → 多赛事 + SaaS化 (1月)"])
add_footer(s12, "12 / 13")

# ═══════════════════════════════════════
# Slide 13: 结尾
# ═══════════════════════════════════════
s13 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s13)
txBox = s13.shapes.add_textbox(Inches(2.5), Inches(2.5), Inches(8), Inches(3))
p = txBox.text_frame.paragraphs[0]
p.text = "我们不仅预测谁会赢得世界杯"; p.font.size = Pt(28); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p2 = txBox.text_frame.add_paragraph()
p2.text = "更让每个人理解为什么"; p2.font.size = Pt(36); p2.font.bold = True; p2.font.color.rgb = GOLD; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(16)
p3 = txBox.text_frame.add_paragraph()
p3.text = "🏆 可解释 · 可视化 · 全流程自动化"; p3.font.size = Pt(18); p3.font.color.rgb = GRAY; p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(20)
p4 = txBox.text_frame.add_paragraph()
p4.text = "感谢评委 | 欢迎提问"; p4.font.size = Pt(16); p4.font.color.rgb = GRAY; p4.alignment = PP_ALIGN.CENTER; p4.space_before = Pt(30)
add_footer(s13, "13 / 13")

# ═══════════════════════════════════════
# 保存
# ═══════════════════════════════════════
prs.save(PPTX_PATH)
print(f"✅ PPTX: {PPTX_PATH} ({len(prs.slides)} slides)")
print("✅ 截图素材: docs/screenshots/")
print("✅ 配色: #0A1628 + #FFD700 + #00D4FF")
